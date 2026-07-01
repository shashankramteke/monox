import collections 
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    prs = Presentation()
    
    # 0 - Title (presentation title slide)
    # 1 - Title and Content
    
    # Slide 1: Title Slide
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "MonoXAI"
    subtitle.text = "High-Precision Forensics & Telemetry Platform\n\nAutomated Instrumentation | Stream Processing | AI Diagnostics"
    for paragraph in title.text_frame.paragraphs:
        paragraph.font.size = Pt(54)
        paragraph.font.bold = True
        paragraph.font.color.rgb = RGBColor(0, 102, 204)
        
    # Slide 2: Introduction
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "What is MonoXAI?"
    tf = content.text_frame
    tf.text = "A premium observability stack designed for multi-service environments."
    p = tf.add_paragraph()
    p.text = "Automated Instrumentation: Effortlessly track services without heavy code changes."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "High-Throughput Stream Processing: Powered by Bytewax for real-time trace reconstruction."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "AI-Powered Diagnostic Engine: Root cause analysis fueled by Google's Gemini."
    p.level = 1

    # Slide 3: System Architecture
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "System Architecture"
    tf = content.text_frame
    tf.text = "End-to-end data flow designed for scale and insights:"
    p = tf.add_paragraph()
    p.text = "Microservices Layer: API Gateway (Node.js) & Quote Service (Python)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Instrumentation: Node & Python Wrappers feeding OTel Collector"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Data Pipeline: OTel Collector -> RabbitMQ -> Bytewax Processor"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Intelligence: FastAPI Backend + SQLite + Gemini AI"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Frontend: React Dashboard (Vite) for Live Analytics"
    p.level = 1
    
    # Slide 4: Forensic Activity Flow
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "Forensic Activity Flow"
    tf = content.text_frame
    tf.text = "How cross-service traces are reconstructed:"
    p = tf.add_paragraph()
    p.text = "User Traffic hits services generating multiple telemetry spans."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "OTel Collector pushes spans, redacts PII, and streams to RabbitMQ."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Bytewax aggregates via Tumbling Window, rebuilding full traces."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Backend requests Forensics -> Gemini AI returns RCA & Fixes."
    p.level = 1

    # Slide 5: Key Features
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "Key Features"
    tf = content.text_frame
    p = tf.paragraphs[0]
    p.text = "AI RCA: Instant root cause analysis with suggested fixes."
    p = tf.add_paragraph()
    p.text = "Trace Waterfall: Multi-service visualization reconstructed in flight."
    p = tf.add_paragraph()
    p.text = "Sparkline Wave & Resource Saturation: Real-time throughput and CPU/Memory."
    p = tf.add_paragraph()
    p.text = "Sidebar Controls: Toggles for Live Mode and Auto-Correlation."
    
    # Slide 6: Understanding Stream Anomalies
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "Understanding Stream Anomalies"
    tf = content.text_frame
    tf.text = "MonoXAI's Stream Processor detects various anomaly streams:"
    p = tf.add_paragraph()
    p.text = "N+1 Query Regression: Detects excessive downstream calls (e.g. DB queries) using Chebyshev bound on span count."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Bimodal Latency: Identifies sudden latency spikes or dual-mode performance using EWMA variance."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Dangling Parent: Finds broken dependency chains (e.g. missing parent spans in distributed traces)."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "PII Redaction Density: Monitors the ratio of PII redactions over a sliding window."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "ML Ensembles: Captures Statistical Outliers, Reconstruction Anomalies, and Boundary Anomalies."
    p.level = 1

    # Slide 7: Viva Questions - Part 1
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "Viva Questions - Part 1"
    tf = content.text_frame
    tf.text = "Q: What is MonoXAI and what problem does it solve?"
    p = tf.add_paragraph()
    p.text = "A: It is a platform for high-precision forensics and telemetry in microservices to easily track distributed applications."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Q: What is an N+1 query and how does MonoXAI detect it?"
    p = tf.add_paragraph()
    p.text = "A: It occurs when an app makes N additional queries to fetch related data. MonoXAI detects it by monitoring span counts using statistical bounds."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Q: How does Bytewax help in trace reconstruction?"
    p = tf.add_paragraph()
    p.text = "A: It aggregates telemetry data streams using time windows (e.g. 10s Tumbling Window) to rebuild scattered cross-service spans into full traces."
    p.level = 1

    # Slide 8: Viva Questions - Part 2
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "Viva Questions - Part 2"
    tf = content.text_frame
    tf.text = "Q: What role does the Gemini AI play in this platform?"
    p = tf.add_paragraph()
    p.text = "A: It performs Root Cause Analysis (RCA) on full traces, identifying errors and suggesting actionable fixes."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Q: Explain Bimodal Latency."
    p = tf.add_paragraph()
    p.text = "A: It happens when a service typically responds quickly but occasionally experiences severe slowdowns, creating two distinct latency modes."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Q: How is PII handled in the pipeline?"
    p = tf.add_paragraph()
    p.text = "A: The OTel Collector redacts PII before streaming. A detector monitors this redaction density to ensure compliance."
    p.level = 1

    # Save presentation
    prs.save("c:/ObserveX-main/MonoXAI_Presentation_Updated.pptx")
    print("Presentation created successfully at c:/ObserveX-main/MonoXAI_Presentation_Updated.pptx")

if __name__ == "__main__":
    create_presentation()
